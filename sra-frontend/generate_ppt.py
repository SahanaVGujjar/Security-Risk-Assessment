"""
Script to generate PowerPoint presentation for Security Risk Assessment Portal
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    primary_color = RGBColor(30, 58, 138)  # #1e3a8a
    accent_color = RGBColor(59, 130, 246)  # #3b82f6
    success_color = RGBColor(34, 197, 94)  # #22c55e
    warning_color = RGBColor(245, 158, 11)  # #f59e0b
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Security Risk Assessment Portal"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Privacy Impact Assessment (PIA) Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    team_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    team_frame = team_box.text_frame
    team_frame.text = "Team Members:\n[Name 1] | [Name 2] | [Name 3]\n\nDate: [Date]"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(18)
    team_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Introduction"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What is a Security Risk Assessment Portal?"
    p = tf.add_paragraph()
    p.text = "• Automated platform for managing Privacy Impact Assessments (PIA) and DPIA"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• System Owners submit detailed assessments about data collection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Approvers review submissions and request clarifications"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automated screening to determine assessment requirements"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Comprehensive audit trail of all responses"
    p.level = 1
    
    # Slide 3: Background
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Background"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Regulatory Landscape"
    p = tf.add_paragraph()
    p.text = "• GDPR requires DPIA for high-risk processing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Organizations must assess privacy risks before processing personal data"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Need for systematic documentation and review processes"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Existing Challenges"
    p = tf.add_paragraph()
    p.text = "• Manual processes are time-consuming"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Lack of standardized questionnaires"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Difficult to track assessment status"
    p.level = 1
    
    # Slide 4: Motivation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Motivation"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Business Drivers"
    p = tf.add_paragraph()
    p.text = "1. Compliance: Meet GDPR/DPIA regulatory requirements"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Efficiency: Automate manual assessment workflows"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Transparency: Maintain clear audit trails"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Collaboration: Enable seamless communication"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Scalability: Handle multiple assessments simultaneously"
    p.level = 1
    
    # Slide 5: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Problems Addressed"
    p = tf.add_paragraph()
    p.text = "Problem 1: Manual Assessment Process"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Email/paper-based submissions, no standardized format"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Problem 2: Lack of Collaboration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• No direct communication, multiple email exchanges"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Problem 3: Inefficient Screening"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Manual determination, no automated routing"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Problem 4: Data Management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Scattered responses, no centralized repository"
    p.level = 2
    
    # Slide 6: Problem Formulation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Formulation"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Functional Requirements"
    p = tf.add_paragraph()
    p.text = "1. User Management: Role-based access (System Owners, Approvers)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Assessment Creation: Create and submit assessments"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Screening Questionnaire: 30+ questions covering data processing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Review Process: Review, request clarifications, approve/reject"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Chat System: Real-time communication for clarifications"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "6. Status Management: Track assessment lifecycle"
    p.level = 1
    
    # Slide 7: Method - Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Method - System Architecture"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technology Stack"
    p = tf.add_paragraph()
    p.text = "Frontend:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• React.js with Material-UI (MUI)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• React Router for navigation"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Context API for authentication"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Backend:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• FastAPI (Python) - RESTful API"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• SQLModel for database ORM"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• SQLite database (scalable to PostgreSQL)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• JWT authentication"
    p.level = 2
    
    # Slide 8: Method - Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Method - Core Features"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System Capabilities"
    p = tf.add_paragraph()
    p.text = "1. Intelligent Screening"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Auto-completion if no personal data collected"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   • Conditional question flow based on answers"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "2. Comprehensive Questionnaire"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • 30+ questions across 5 pages"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   • Multiple question types: Boolean, Text, Multiselect, Dropdown"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "3. Collaboration Tools"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Thread-based chat system"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "   • Multiple clarifications per question"
    p.level = 2
    
    # Slide 9: Method - Workflows
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Method - User Workflows"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System Owner Workflow"
    p = tf.add_paragraph()
    p.text = "1. Create new assessment"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Answer screening questions (paginated)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Submit assessment"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Receive clarification requests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "5. Edit responses and resubmit"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Approver Workflow"
    p = tf.add_paragraph()
    p.text = "1. View all submitted assessments"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Review responses and raise clarifications"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Mark as Completed or Red Flag"
    p.level = 1
    
    # Slide 10: Experiments - Demo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Preliminary Experiments - System Demo"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Functionalities Demonstrated"
    p = tf.add_paragraph()
    p.text = "Demo Flow:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "1. Login & Dashboard: Role-based access"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "2. Create Assessment: System Owner creates new assessment"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "3. Screening Process: Answer 30+ questions with conditional logic"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "4. Review Process: Approver views assessment"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "5. Clarification: Approver raises questions on multiple questions"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "6. Response: System Owner edits and responds"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "7. Completion: Approver marks as completed"
    p.level = 2
    
    # Slide 11: Experiments - Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Preliminary Experiments - Results"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System Performance & Usability"
    p = tf.add_paragraph()
    p.text = "Performance Metrics:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Fast page load times (< 2 seconds)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Smooth pagination between question pages"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Efficient API responses"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "User Experience:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Intuitive navigation"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Clear visual indicators (orange borders for clarifications)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Responsive design (works on mobile/tablet)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Challenges Addressed: ✅ Automated screening ✅ Standardized questionnaire ✅ Centralized repository ✅ Real-time collaboration"
    p.level = 1
    
    # Slide 12: Conclusions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Achievements"
    p = tf.add_paragraph()
    p.text = "✅ Fully functional SRA Portal"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Role-based access control"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Comprehensive screening questionnaire"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Real-time collaboration features"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Complete audit trail"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Enhancements"
    p = tf.add_paragraph()
    p.text = "1. Advanced Analytics: Dashboard with assessment statistics"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Export Functionality: PDF/Excel reports"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Email Notifications: Automated alerts"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Multi-tenant Support: Organization-level isolation"
    p.level = 1
    
    # Slide 13: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    q_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    q_frame = q_box.text_frame
    q_frame.text = "Questions?"
    q_para = q_frame.paragraphs[0]
    q_para.font.size = Pt(32)
    q_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Contact: [Email addresses]\nGitHub: [Repository URL]"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(18)
    contact_para.alignment = PP_ALIGN.CENTER
    
    prs.save('Security_Risk_Assessment_Portal_Presentation.pptx')
    print("Presentation created successfully: Security_Risk_Assessment_Portal_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx not installed. Please run: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentation: {e}")

